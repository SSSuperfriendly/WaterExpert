from __future__ import annotations

import argparse
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches


CANVAS = (3600, 1600)
FIG_LEFT = Inches(0.48)
FIG_TOP = Inches(1.08)
FIG_WIDTH = Inches(12.25)


def load_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        r"C:\Windows\Fonts\msyhbd.ttc" if bold else r"C:\Windows\Fonts\msyh.ttc",
        r"C:\Windows\Fonts\simhei.ttf" if bold else r"C:\Windows\Fonts\simsun.ttc",
    ]
    for path in candidates:
        p = Path(path)
        if p.exists():
            return ImageFont.truetype(str(p), size=size)
    return ImageFont.load_default()


def rounded_box(draw: ImageDraw.ImageDraw, box, fill, outline=None, radius=26, shadow=True, shadow_color=(23, 59, 122, 28)):
    x0, y0, x1, y1 = box
    if shadow:
        draw.rounded_rectangle((x0 + 8, y0 + 10, x1 + 8, y1 + 10), radius=radius, fill=shadow_color)
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=3 if outline else 0)


def gradient_background(size, top_color=(248, 251, 255), bottom_color=(255, 255, 255)):
    w, h = size
    img = Image.new("RGBA", size, (255, 255, 255, 255))
    px = img.load()
    for y in range(h):
        t = y / max(h - 1, 1)
        color = tuple(int(top_color[i] * (1 - t) + bottom_color[i] * t) for i in range(3)) + (255,)
        for x in range(w):
            px[x, y] = color
    return img


def add_decor(img: Image.Image) -> None:
    draw = ImageDraw.Draw(img, "RGBA")
    w, h = img.size
    draw.ellipse((-180, -160, 620, 500), fill=(58, 114, 222, 24))
    draw.ellipse((w - 700, -120, w + 80, 520), fill=(63, 109, 179, 18))
    draw.ellipse((w - 560, h - 280, w + 140, h + 180), fill=(44, 160, 164, 20))

    for x in range(120, w - 120, 120):
        draw.line((x, 120, x, h - 120), fill=(190, 206, 230, 32), width=1)
    for y in range(160, h - 120, 120):
        draw.line((110, y, w - 110, y), fill=(190, 206, 230, 28), width=1)


def draw_stage_header(draw, box, number, title, palette):
    x0, y0, x1, _ = box
    pill_h = 78
    rounded_box(draw, (x0, y0, x1, y0 + pill_h), fill=palette["deep"], radius=28, shadow=False)
    circle = (x0 + 22, y0 + 15, x0 + 70, y0 + 63)
    draw.ellipse(circle, fill=(255, 255, 255, 255))
    num_font = load_font(28, bold=True)
    title_font = load_font(34, bold=True)
    nbox = draw.textbbox((0, 0), number, font=num_font)
    draw.text((circle[0] + 24 - (nbox[2] - nbox[0]) / 2, y0 + 23), number, font=num_font, fill=palette["deep"])
    draw.text((x0 + 88, y0 + 18), title, font=title_font, fill=(255, 255, 255, 255))


def draw_stage_panel(draw, box, palette):
    x0, y0, x1, y1 = box
    rounded_box(draw, box, fill=palette["panel"], outline=palette["outline"], radius=36)
    draw.line((x0, y0 + 82, x1, y0 + 82), fill=palette["outline"], width=2)


def draw_card(draw, box, title, body_lines, accent, palette, body_size=21):
    x0, y0, x1, y1 = box
    rounded_box(draw, box, fill=(255, 255, 255, 255), outline=(214, 224, 240, 255), radius=24, shadow=True)
    draw.rounded_rectangle((x0, y0, x1, y0 + 10), radius=24, fill=accent)
    draw.rectangle((x0, y0 + 6, x1, y0 + 10), fill=accent)
    tfont = load_font(26, bold=True)
    bfont = load_font(body_size, bold=False)
    draw.text((x0 + 26, y0 + 26), title, font=tfont, fill=palette["title"])
    y = y0 + 74
    for line in body_lines:
        draw.text((x0 + 28, y), f"• {line}", font=bfont, fill=palette["body"])
        y += body_size + 12


def draw_arrow(draw, start, end, color, width=10):
    x0, y0 = start
    x1, y1 = end
    draw.line((x0, y0, x1, y1), fill=color, width=width)
    head = 24
    draw.polygon([(x1, y1), (x1 - head, y1 - head // 2), (x1 - head, y1 + head // 2)], fill=color)


def draw_chips(draw, origin, labels, fill, text_fill, chip_gap=18):
    x, y = origin
    font = load_font(24, bold=True)
    for label in labels:
        bbox = draw.textbbox((0, 0), label, font=font)
        w = bbox[2] - bbox[0] + 40
        h = 48
        rounded_box(draw, (x, y, x + w, y + h), fill=fill, radius=22, shadow=False)
        draw.text((x + 20, y + 11), label, font=font, fill=text_fill)
        x += w + chip_gap


def mscim_palette():
    return {
        "deep": (34, 84, 174, 255),
        "panel": (243, 248, 255, 230),
        "outline": (161, 192, 239, 255),
        "title": (23, 59, 122, 255),
        "body": (72, 91, 118, 255),
        "teal": (58, 165, 170, 255),
        "orange": (241, 130, 45, 255),
        "violet": (106, 127, 196, 255),
        "green": (47, 143, 92, 255),
    }


def cmfbe_palette():
    return {
        "deep": (22, 74, 131, 255),
        "panel": (246, 250, 255, 232),
        "outline": (173, 204, 227, 255),
        "title": (18, 53, 106, 255),
        "body": (73, 92, 115, 255),
        "warm1": (212, 80, 99, 255),
        "warm2": (246, 151, 74, 255),
        "cool1": (62, 166, 167, 255),
        "cool2": (94, 150, 140, 255),
        "neutral": (91, 115, 168, 255),
        "green": (45, 142, 92, 255),
    }


def make_mscim_figure(output_path: Path):
    palette = mscim_palette()
    img = gradient_background(CANVAS)
    add_decor(img)
    draw = ImageDraw.Draw(img, "RGBA")

    cols = [
        (90, 110, 760, 1125),
        (860, 110, 1530, 1125),
        (1630, 110, 2640, 1125),
        (2740, 110, 3500, 1125),
    ]
    titles = [("1", "多源输入层"), ("2", "统一对齐与构图"), ("3", "MSCIM 核心建模"), ("4", "输出与应用")]

    for box, (n, t) in zip(cols, titles):
        draw_stage_panel(draw, box, palette)
        draw_stage_header(draw, (box[0], box[1], box[1] + 1 if False else box[2], box[3]), n, t, palette)

    # Stage 1
    stage1_cards = [
        ("水质监测", ["温度 / pH / DO / 电导率", "浊度 / TN / TP / 氨氮等"], palette["deep"]),
        ("天气驱动", ["气压 / 气温 / 湿度 / 降水", "风速 / 风向"], (78, 150, 224, 255)),
        ("水动力背景", ["水位 / 流量 / 潮汐影响", "当前参考：松浦大桥、黄渡"], palette["teal"]),
        ("文本与先验", ["工程案例 / 治理经验", "机理规则 / 站点属性"], palette["violet"]),
    ]
    y = 214
    for title, lines, accent in stage1_cards:
        draw_card(draw, (118, y, 732, y + 176), title, lines, accent, palette)
        y += 198

    # Stage 2
    stage2_cards = [
        ("站点-日期对齐", ["统一主键：station × date", "形成可追踪日尺度样本"], palette["deep"]),
        ("质量控制", ["异常值处理 / 缺失控制", "透明度 proxy 与天气匹配"], (83, 150, 222, 255)),
        ("时空图构建", ["时间窗口序列", "站点邻接 + 因果先验矩阵"], (101, 125, 197, 255)),
        ("知识增强特征", ["机理规则与案例知识", "编码为结构化输入特征"], palette["teal"]),
    ]
    y = 214
    for title, lines, accent in stage2_cards:
        draw_card(draw, (888, y, 1502, y + 176), title, lines, accent, palette)
        y += 198

    # Stage 3
    model_boxes = [
        ((1668, 206, 2602, 388), "Transformer 时序编码", ["学习季节性、滞后效应", "捕捉突变与长期依赖"], palette["deep"]),
        ((1668, 420, 2602, 602), "时空因果融合", ["站点关系 + 因果注意力", "识别主导因子贡献强度"], palette["teal"]),
        ((1668, 634, 2602, 864), "多任务联合头", ["预测：清澈度 / 浊度", "识别：重点治理区边界", "诊断：致浊因子排序"], palette["orange"]),
    ]
    for box, title, lines, accent in model_boxes:
        draw_card(draw, box, title, lines, accent, palette, body_size=22)

    center_box = (1690, 908, 2580, 1068)
    rounded_box(draw, center_box, fill=(236, 244, 255, 255), outline=(198, 214, 240, 255), radius=28)
    draw.text((1722, 936), "模型定位", font=load_font(26, bold=True), fill=palette["title"])
    draw.text((1722, 984), "面向“预测 + 诊断”一体化，不只给出数值结果，", font=load_font(22), fill=palette["body"])
    draw.text((1722, 1018), "还输出重点治理区边界与主导致浊因子。", font=load_font(22), fill=palette["body"])

    # Stage 4
    stage4_cards = [
        ("动态预测", ["清澈度 / 浊度日尺度变化", "支撑趋势预警"], palette["orange"]),
        ("边界识别", ["重点治理区 - 稳定区", "空间边界自动识别"], (244, 163, 92, 255)),
        ("致因诊断", ["给出主导致浊因子", "形成可解释因果链"], (218, 117, 33, 255)),
        ("治理支撑", ["支撑调度与优先级排序", "连接情景分析设计"], (203, 103, 30, 255)),
    ]
    y = 214
    for title, lines, accent in stage4_cards:
        draw_card(draw, (2772, y, 3468, y + 176), title, lines, accent, palette)
        y += 198

    ay = 618
    draw_arrow(draw, (782, ay), (838, ay), palette["deep"], width=11)
    draw_arrow(draw, (1552, ay), (1608, ay), palette["deep"], width=11)
    draw_arrow(draw, (2662, ay), (2718, ay), palette["deep"], width=11)

    # Bottom callout
    bottom = (90, 1190, 3500, 1510)
    rounded_box(draw, bottom, fill=(241, 250, 244, 255), outline=(165, 215, 184, 255), radius=34)
    draw.text((122, 1230), "当前落地情况", font=load_font(32, bold=True), fill=palette["green"])
    draw_chips(
        draw,
        (360, 1224),
        [
            "全站基础库 20 个主库站点",
            "31099 条日尺度记录",
            "水质 + 透明度 proxy + 天气",
        ],
        fill=(225, 243, 232, 255),
        text_fill=(42, 109, 72, 255),
    )
    draw_chips(
        draw,
        (360, 1292),
        [
            "单站增强对象：吴淞口",
            "主参考：松浦大桥",
            "辅助参考：黄渡",
            "可训练重叠：891 天",
        ],
        fill=(228, 238, 255, 255),
        text_fill=(41, 76, 140, 255),
    )
    draw.text((360, 1388), "下一步：补遥感 / NDTI / 治理事件 / 更细颗粒度断面水动力，继续强化时空因果诊断能力。", font=load_font(25), fill=(51, 94, 66, 255))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(output_path)


def make_cmfbe_figure(output_path: Path):
    palette = cmfbe_palette()
    img = gradient_background(CANVAS, top_color=(248, 251, 255), bottom_color=(255, 255, 255))
    add_decor(img)
    draw = ImageDraw.Draw(img, "RGBA")

    cols = [
        (90, 110, 760, 1125),
        (860, 110, 1570, 1125),
        (1670, 110, 2650, 1125),
        (2750, 110, 3500, 1125),
    ]
    titles = [("1", "驱动与边界"), ("2", "过程项分解"), ("3", "CMFBE-ST-GCN 核心"), ("4", "输出与应用")]
    for box, (n, t) in zip(cols, titles):
        draw_stage_panel(draw, box, palette)
        draw_stage_header(draw, (box[0], box[1], box[2], box[3]), n, t, palette)

    # Stage 1
    stage1_cards = [
        ("外部驱动", ["降雨径流 / 风浪扰动", "气温 / 湿度 / 辐射背景"], palette["neutral"]),
        ("水动力边界", ["水位 / 流量 / 回水效应", "潮汐滞留 / 冲刷外输"], palette["cool1"]),
        ("水质状态", ["浊度 / 透明度 / 电导率", "营养盐 / 藻类扰动"], palette["warm1"]),
        ("调度与情景", ["治理事件 / 工程调度", "情景参数与边界条件"], palette["green"]),
    ]
    y = 214
    for title, lines, accent in stage1_cards:
        draw_card(draw, (118, y, 732, y + 176), title, lines, accent, palette)
        y += 198

    # Stage 2: warm/cool process decomposition
    source_box = (888, 214, 1542, 540)
    sink_box = (888, 590, 1542, 916)
    rounded_box(draw, source_box, fill=(255, 245, 245, 255), outline=(244, 189, 199, 255), radius=28)
    rounded_box(draw, sink_box, fill=(241, 251, 251, 255), outline=(169, 219, 220, 255), radius=28)
    draw.text((920, 244), "致浊源项（Source terms）", font=load_font(30, bold=True), fill=palette["warm1"])
    draw_chips(draw, (920, 298), ["径流输入", "潮汐滞留", "再悬浮", "生物扰动"], fill=(253, 226, 231, 255), text_fill=(148, 48, 72, 255), chip_gap=14)
    draw.text((920, 378), "表征使水体变浑的输入、滞留与扰动过程。", font=load_font(22), fill=palette["body"])
    draw.text((920, 620), "去浊汇项（Sink terms）", font=load_font(30, bold=True), fill=palette["cool1"])
    draw_chips(draw, (920, 674), ["冲刷外输", "沉降絮凝", "自净恢复", "稀释交换"], fill=(221, 244, 244, 255), text_fill=(34, 122, 122, 255), chip_gap=14)
    draw.text((920, 754), "表征使水体恢复变清的输出、沉降与恢复过程。", font=load_font(22), fill=palette["body"])

    transport_box = (908, 962, 1520, 1072)
    rounded_box(draw, transport_box, fill=(245, 248, 255, 255), outline=(207, 217, 239, 255), radius=24)
    draw.text((936, 994), "输移重分配：平流-扩散 / 回水-滞留 / 空间传播", font=load_font(24, bold=True), fill=palette["neutral"])

    # Stage 3
    core1 = (1702, 206, 2618, 380)
    core2 = (1702, 406, 2618, 602)
    core3 = (1702, 630, 2618, 870)
    core4 = (1702, 900, 2618, 1070)
    draw_card(draw, core1, "过程分支网络", ["分别学习源项、汇项和输移项", "避免把物理过程混成单一黑箱"], palette["neutral"], palette)
    draw_card(draw, core2, "ST-GCN 图传播", ["传播站点间相互影响", "显式表达空间耦合与时间依赖"], palette["deep"], palette)
    draw_card(draw, core3, "机理约束损失", ["质量守恒 / 符号一致 / 阈值响应", "提升情景泛化与可解释性"], palette["green"], palette)
    rounded_box(draw, core4, fill=(239, 246, 255, 255), outline=(204, 220, 242, 255), radius=24)
    draw.text((1732, 932), "净变化表达", font=load_font(28, bold=True), fill=palette["title"])
    draw.text((1732, 984), "dC/dt ≈ Source terms  −  Sink terms  +  Transport redistribution", font=load_font(24, bold=True), fill=palette["neutral"])

    # Stage 4
    stage4_cards = [
        ("动态模拟", ["浊度 / 清澈度轨迹", "恢复过程连续模拟"], palette["warm2"]),
        ("过程归因", ["解释转折点由哪些过程主导", "形成源-汇分解证据"], palette["warm1"]),
        ("阈值瓶颈", ["识别限制清澈度提升的关键阈值", "锁定瓶颈环节"], palette["cool1"]),
        ("情景比较", ["调度 / 工程 / 外部扰动对比", "支撑治理策略设计"], palette["green"]),
    ]
    y = 214
    for title, lines, accent in stage4_cards:
        draw_card(draw, (2782, y, 3468, y + 176), title, lines, accent, palette)
        y += 198

    ay = 618
    draw_arrow(draw, (782, ay), (838, ay), palette["deep"], width=11)
    draw_arrow(draw, (1592, ay), (1648, ay), palette["deep"], width=11)
    draw_arrow(draw, (2672, ay), (2728, ay), palette["deep"], width=11)

    # bottom strip
    bottom = (90, 1190, 3500, 1510)
    rounded_box(draw, bottom, fill=(247, 249, 253, 255), outline=(194, 210, 232, 255), radius=34)
    draw.text((122, 1228), "当前实现进展", font=load_font(32, bold=True), fill=palette["deep"])
    draw_chips(
        draw,
        (360, 1222),
        ["已实现致浊-去浊过程分解", "可解释测试阶段变浑/恢复机制", "与 MSCIM 形成“诊断 + 机理模拟”互补"],
        fill=(232, 239, 252, 255),
        text_fill=(36, 77, 140, 255),
        chip_gap=14,
    )
    draw_chips(
        draw,
        (360, 1290),
        ["Warm = 致浊源项", "Cool = 去浊汇项", "Physics-informed = 约束模型不违背基本机理"],
        fill=(240, 248, 248, 255),
        text_fill=(46, 107, 109, 255),
        chip_gap=14,
    )
    draw.text((360, 1388), "CMFBE-ST-GCN 的重点不是单纯提高拟合，而是把“为什么变浑、为什么恢复变清”拆成可解释过程，并用于情景模拟。", font=load_font(24), fill=(70, 86, 112, 255))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(output_path)


def replace_or_add_picture(slide, name_prefix: str, image_path: Path):
    removable = []
    for shape in slide.shapes:
        name = getattr(shape, "name", "")
        if name.startswith(name_prefix):
            removable.append(shape)
    for shape in removable:
        sp = shape._element
        sp.getparent().remove(sp)
    pic = slide.shapes.add_picture(str(image_path), FIG_LEFT, FIG_TOP, width=FIG_WIDTH)
    pic.name = f"{name_prefix}_MAIN"


def cleanup_slide_content(slide, remove_extra_text=False):
    removable = []
    for shape in slide.shapes:
        st = shape.shape_type
        text = ""
        if hasattr(shape, "text") and shape.text:
            text = shape.text.strip()
        if st == 13 and shape.top > 900000:  # picture in content area
            removable.append(shape)
        elif remove_extra_text and st == 17 and shape.top > 1000000 and "2.4" not in text and "2.3" not in text:
            removable.append(shape)
    for shape in removable:
        sp = shape._element
        sp.getparent().remove(sp)


def build_ppt(source_ppt: Path, output_ppt: Path, mscim_image: Path, cmfbe_image: Path):
    prs = Presentation(str(source_ppt))
    cleanup_slide_content(prs.slides[5], remove_extra_text=False)
    cleanup_slide_content(prs.slides[6], remove_extra_text=True)
    replace_or_add_picture(prs.slides[5], "MSCIM_POLISHED_DIAGRAM", mscim_image)
    replace_or_add_picture(prs.slides[6], "CMFBE_POLISHED_DIAGRAM", cmfbe_image)
    if output_ppt.exists():
        output_ppt.unlink()
    prs.save(str(output_ppt))


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--ppt", type=Path, required=True)
    parser.add_argument("--output-ppt", type=Path, required=True)
    parser.add_argument("--mscim-image", type=Path, required=True)
    parser.add_argument("--cmfbe-image", type=Path, required=True)
    args = parser.parse_args()

    make_mscim_figure(args.mscim_image)
    make_cmfbe_figure(args.cmfbe_image)
    build_ppt(args.ppt, args.output_ppt, args.mscim_image, args.cmfbe_image)
    print(f"mscim_image={args.mscim_image}")
    print(f"cmfbe_image={args.cmfbe_image}")
    print(f"ppt_saved={args.output_ppt}")


if __name__ == "__main__":
    main()
