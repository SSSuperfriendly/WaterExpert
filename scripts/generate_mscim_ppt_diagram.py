from __future__ import annotations

import argparse
import shutil
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches


def load_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        r"C:\Windows\Fonts\msyhbd.ttc" if bold else r"C:\Windows\Fonts\msyh.ttc",
        r"C:\Windows\Fonts\simhei.ttf" if bold else r"C:\Windows\Fonts\simsun.ttc",
    ]
    for path in candidates:
        if Path(path).exists():
            return ImageFont.truetype(path, size=size)
    return ImageFont.load_default()


def draw_shadowed_round_box(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    fill: str,
    outline: str,
    radius: int = 28,
    shadow_offset: int = 8,
) -> None:
    x0, y0, x1, y1 = box
    shadow = (x0 + shadow_offset, y0 + shadow_offset, x1 + shadow_offset, y1 + shadow_offset)
    draw.rounded_rectangle(shadow, radius=radius, fill="#D9E2F3")
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=3)


def draw_panel(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    number: str,
    title: str,
    fill: str,
    header_fill: str = "#1F4EAA",
) -> None:
    x0, y0, x1, y1 = box
    draw_shadowed_round_box(draw, box, fill=fill, outline="#A9BFE5", radius=34)
    header_h = 96
    draw.rounded_rectangle((x0, y0, x1, y0 + header_h), radius=34, fill=header_fill)
    draw.rectangle((x0, y0 + 44, x1, y0 + header_h), fill=header_fill)

    num_r = 34
    cx = x0 + 52
    cy = y0 + header_h // 2
    draw.ellipse((cx - num_r, cy - num_r, cx + num_r, cy + num_r), fill="#FFFFFF")

    font_num = load_font(34, bold=True)
    font_title = load_font(34, bold=True)
    num_bbox = draw.textbbox((0, 0), number, font=font_num)
    draw.text((cx - (num_bbox[2] - num_bbox[0]) / 2, cy - (num_bbox[3] - num_bbox[1]) / 2 - 2), number, font=font_num, fill=header_fill)
    draw.text((x0 + 102, y0 + 24), title, font=font_title, fill="#FFFFFF")


def draw_card(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    title: str,
    body_lines: list[str],
    accent: str,
) -> None:
    x0, y0, x1, y1 = box
    draw_shadowed_round_box(draw, box, fill="#FFFFFF", outline="#D5DFF0", radius=24, shadow_offset=6)
    draw.rounded_rectangle((x0, y0, x1, y0 + 14), radius=24, fill=accent)
    draw.rectangle((x0, y0 + 8, x1, y0 + 14), fill=accent)

    title_font = load_font(25, bold=True)
    body_font = load_font(21, bold=False)
    draw.text((x0 + 24, y0 + 28), title, font=title_font, fill="#173B7A")

    line_y = y0 + 76
    for line in body_lines:
        draw.text((x0 + 26, line_y), f"• {line}", font=body_font, fill="#334E68")
        line_y += 34


def draw_arrow(draw: ImageDraw.ImageDraw, start: tuple[int, int], end: tuple[int, int], color: str = "#1F4EAA") -> None:
    x0, y0 = start
    x1, y1 = end
    draw.line((x0, y0, x1, y1), fill=color, width=10)
    head = 24
    draw.polygon(
        [
            (x1, y1),
            (x1 - head, y1 - head // 2),
            (x1 - head, y1 + head // 2),
        ],
        fill=color,
    )


def build_diagram(output_image: Path) -> None:
    width, height = 3400, 1520
    img = Image.new("RGB", (width, height), "#FFFFFF")
    draw = ImageDraw.Draw(img)

    panel_y0 = 84
    panel_y1 = 1120
    p1 = (90, panel_y0, 790, panel_y1)
    p2 = (910, panel_y0, 1590, panel_y1)
    p3 = (1710, panel_y0, 2640, panel_y1)
    p4 = (2760, panel_y0, 3320, panel_y1)

    draw_panel(draw, p1, "1", "多源输入层", fill="#EDF4FF")
    draw_panel(draw, p2, "2", "统一对齐与构图", fill="#F5F8FE")
    draw_panel(draw, p3, "3", "MSCIM 核心建模", fill="#EEF6FF")
    draw_panel(draw, p4, "4", "输出与应用", fill="#FFF7ED")

    cards1 = [
        ("水质监测", ["水温 / pH / DO / 电导率", "浊度 / TN / TP / 氨氮等"], "#4F81BD"),
        ("天气驱动", ["气压 / 气温 / 湿度", "降水 / 风速 / 风向"], "#5B9BD5"),
        ("水动力背景", ["水位 / 流量", "当前参考站：松浦大桥、黄渡"], "#4CA6A8"),
        ("文本与先验", ["工程案例 / 治理经验", "机理规则 / 站点属性"], "#7E93C8"),
    ]
    y = 206
    for title, lines, accent in cards1:
        draw_card(draw, (122, y, 758, y + 184), title, lines, accent)
        y += 212

    cards2 = [
        ("站点-日期对齐", ["统一主键：station × date", "形成可追踪日尺度样本"], "#4F81BD"),
        ("质量控制", ["异常值处理 / 缺失控制", "透明度 proxy 与天气匹配"], "#5B9BD5"),
        ("时空图构建", ["时间窗口序列", "站点邻接 + 因果先验矩阵"], "#7E93C8"),
        ("知识增强特征", ["把机理规则与案例知识", "编码为结构化输入特征"], "#4CA6A8"),
    ]
    y = 206
    for title, lines, accent in cards2:
        draw_card(draw, (944, y, 1558, y + 184), title, lines, accent)
        y += 212

    cards3 = [
        ("Transformer 时序编码", ["学习季节性、滞后效应", "捕捉突变与长期依赖"], "#3F6DB3"),
        ("时空因果融合", ["站点关系 + 因果注意力", "识别主导因子贡献强度"], "#4CA6A8"),
        ("多任务联合头", ["预测头：清澈度 / 浊度", "识别头：重点治理区边界", "诊断头：致浊因子排序"], "#ED7D31"),
    ]
    y = 190
    heights = [206, 206, 232]
    for idx, (title, lines, accent) in enumerate(cards3):
        h = heights[idx]
        draw_card(draw, (1746, y, 2606, y + h), title, lines, accent)
        y += h + 28

    pseudo_box = (1760, 906, 2592, 1080)
    draw_shadowed_round_box(draw, pseudo_box, fill="#173B7A", outline="#173B7A", radius=24, shadow_offset=6)
    pseudo_title = load_font(24, bold=True)
    pseudo_body = load_font(22, bold=False)
    draw.text((1788, 930), "伪算法表达", font=pseudo_title, fill="#FFFFFF")
    pseudo_lines = [
        "for station s, date t:",
        "  X[s,t] = Align(水质, 天气, 水动力, 知识)",
        "  H[t] = Transformer(X[s,t-L:t])",
        "  Z[t] = CausalFusion(H[t], G[t], K)",
        "  输出 = 预测 + 边界识别 + 致因诊断",
    ]
    py = 968
    for line in pseudo_lines:
        draw.text((1790, py), line, font=pseudo_body, fill="#E9F2FF")
        py += 24

    cards4 = [
        ("动态预测", ["清澈度 / 浊度的", "日尺度变化趋势"], "#ED7D31"),
        ("边界识别", ["重点治理区 - 稳定区", "空间边界自动识别"], "#F29F67"),
        ("致因诊断", ["给出主导致浊因子", "形成可解释因果链"], "#D97A28"),
        ("治理支撑", ["支撑调度、治理优先级", "与情景分析设计"], "#C96A1B"),
    ]
    y = 206
    for title, lines, accent in cards4:
        draw_card(draw, (2794, y, 3288, y + 184), title, lines, accent)
        y += 212

    arrow_y = 610
    draw_arrow(draw, (806, arrow_y), (892, arrow_y))
    draw_arrow(draw, (1608, arrow_y), (1690, arrow_y))
    draw_arrow(draw, (2656, arrow_y), (2740, arrow_y))

    bottom_box = (90, 1180, 3320, 1440)
    draw_shadowed_round_box(draw, bottom_box, fill="#EFFAF2", outline="#A8D5B7", radius=30, shadow_offset=8)
    tag_box = (122, 1214, 392, 1290)
    draw.rounded_rectangle(tag_box, radius=20, fill="#2E8B57")
    tag_font = load_font(28, bold=True)
    body_font = load_font(25, bold=False)
    body_bold = load_font(26, bold=True)
    draw.text((152, 1233), "当前落地情况", font=tag_font, fill="#FFFFFF")
    draw.text((432, 1218), "全站基础库：20 个主库站点 / 31099 条日尺度记录 / 水质 + 透明度 proxy + 天气", font=body_bold, fill="#1F5130")
    draw.text((432, 1270), "单站增强：吴淞口为目标水质站，松浦大桥为主参考水动力站，黄渡为辅助参考站", font=body_font, fill="#2B5D3C")
    draw.text((432, 1320), "当前可直接用于水动力增强训练的重叠时段：2022-01-01 至 2024-12-31，共 891 天", font=body_font, fill="#2B5D3C")
    draw.text((432, 1370), "下一步重点：补遥感 / NDTI / 治理事件 / 更细颗粒度断面水动力，实现更强的时空因果诊断", font=body_font, fill="#2B5D3C")

    footer_font = load_font(20, bold=False)
    draw.text((2460, 1464), "MSCIM = 多源输入 × 时空编码 × 因果解释 × 治理决策支撑", font=footer_font, fill="#5B6B82")

    output_image.parent.mkdir(parents=True, exist_ok=True)
    img.save(output_image, quality=95)


def inject_into_ppt(source_ppt: Path, output_ppt: Path, image_path: Path, slide_index: int = 6) -> None:
    prs = Presentation(str(source_ppt))
    slide = prs.slides[slide_index - 1]

    # Remove previously inserted image if rerun.
    removable = []
    for shape in slide.shapes:
        if getattr(shape, "name", "").startswith("MSCIM_LOGIC_DIAGRAM"):
            removable.append(shape)
    for shape in removable:
        sp = shape._element
        sp.getparent().remove(sp)

    pic = slide.shapes.add_picture(str(image_path), Inches(0.55), Inches(1.15), width=Inches(12.2))
    pic.name = "MSCIM_LOGIC_DIAGRAM_MAIN"
    if output_ppt.exists():
        output_ppt.unlink()
    prs.save(str(output_ppt))


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate MSCIM logic diagram and inject it into a PPT.")
    parser.add_argument("--ppt", type=Path, required=True, help="Source PPTX path.")
    parser.add_argument("--output-ppt", type=Path, required=True, help="Output PPTX path.")
    parser.add_argument("--output-image", type=Path, required=True, help="Output image path.")
    args = parser.parse_args()

    build_diagram(args.output_image)
    inject_into_ppt(args.ppt, args.output_ppt, args.output_image)

    print(f"image_saved={args.output_image}")
    print(f"ppt_saved={args.output_ppt}")


if __name__ == "__main__":
    main()
